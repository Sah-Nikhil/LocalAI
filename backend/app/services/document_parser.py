# app/services/document_parser.py

import os
import re
import base64
import requests
from io import BytesIO
from PIL import Image
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from app.lib.load_config import get_vlm_model_name, get_vlm_host, get_model_name, get_llm_host

# Regex to match figure captions like "Fig 1.1 - Title"
FIGURE_CAPTION_REGEX = re.compile(r"(Fig(?:ure)?\.?\s*\d+(?:\.\d+)*)(?:[:\-\u2013]\s*(.*))?", re.IGNORECASE)

def describe_image_with_qwen(image: Image.Image, caption: str = "") -> str:
    buffered = BytesIO()
    image.save(buffered, format="PNG")
    img_b64 = base64.b64encode(buffered.getvalue()).decode("utf-8")

    prompt = "Describe this image in detail, including any diagrams, arrows, labels, or flows."
    if caption:
        prompt += f" This image is captioned: '{caption}'. Please use that information to improve your interpretation."

    payload = {
        "model": get_vlm_model_name(),
        "prompt": prompt,
        "images": [img_b64],
        "stream": False
    }

    try:
        response = requests.post(f"{get_vlm_host()}/api/generate", json=payload, timeout=30)
        result = response.json()
        return result.get("response", "[No description generated]").strip()
    except Exception as e:
        return f"[VLM Error: {str(e)}]"

def extract_text(path: str) -> str:
    ext = os.path.splitext(path)[-1].lower()

    if ext == ".pdf":
        combined_chunks = []
        doc = fitz.open(path)

        for i, page in enumerate(doc):
            page_text = page.get_text("text").strip()
            page_summary = [f"[Page {i+1}]\n{page_text}"]

            image_list = page.get_images(full=False)
            for img_index, img in enumerate(image_list):
                try:
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    img = Image.open(BytesIO(image_bytes))

                    caption_matches = FIGURE_CAPTION_REGEX.findall(page_text)
                    caption = caption_matches[-1][1] if caption_matches else ""

                    print(f"🖼️ Interpreting image on PDF page {i+1}, image {img_index+1} via VLM")
                    desc = describe_image_with_qwen(img, caption)
                    page_summary.append(f"[Image Summary] {desc}")
                except Exception as e:
                    print(f"⚠️ Failed to process image on PDF page {i+1}, image {img_index+1}: {e}")

            combined_chunks.append("\n".join(page_summary))

        return "\n\n".join(combined_chunks)

    elif ext == ".docx":
        doc = Document(path)
        text_parts = [p.text for p in doc.paragraphs if p.text.strip()]

        last_caption = ""
        for rel in doc.part._rels:
            rel_obj = doc.part._rels[rel]
            if "image" in rel_obj.target_ref:
                img_data = rel_obj.target_part.blob
                img = Image.open(BytesIO(img_data))

                caption_match = next((m for m in reversed(text_parts) if FIGURE_CAPTION_REGEX.search(m)), None)
                if caption_match:
                    match = FIGURE_CAPTION_REGEX.search(caption_match)
                    last_caption = match.group(2) if match else ""

                print("🖼️ Interpreting embedded Word image via VLM")
                desc = describe_image_with_qwen(img, last_caption)
                text_parts.append(f"[VLM DOCX Image]\n{desc}")

        return "\n\n".join(text_parts)

    elif ext == ".pptx":
        prs = Presentation(path)
        all_text = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text = [f"[Slide {slide_num}]"]
            possible_captions = []

            for shape in slide.shapes:
                if hasattr(shape, "text") and isinstance(getattr(shape, "text", None), str) and shape.text.strip():
                    text = shape.text.strip()
                    slide_text.append(text)
                    if FIGURE_CAPTION_REGEX.search(text):
                        possible_captions.append(text)

                if hasattr(shape, "has_table") and shape.has_table and hasattr(shape, "table"):
                    for row in shape.table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells)
                        slide_text.append(row_text)

                if hasattr(shape, "image") and getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image = shape.image
                        img_bytes = image.blob
                        img = Image.open(BytesIO(img_bytes))
                        caption = possible_captions[-1] if possible_captions else ""
                        print(f"🖼️ Interpreting image on Slide {slide_num} via VLM")
                        desc = describe_image_with_qwen(img, caption)
                        slide_text.append(f"[VLM Slide Image]\n{desc}")
                    except Exception:
                        pass

            all_text.append("\n".join(slide_text))

        return "\n\n".join(all_text)

    elif ext in [".txt", ".md"]:
        with open(path, "r", encoding="utf-8") as f:
            return f.read()

    else:
        raise ValueError(f"Unsupported file type: {ext}")

def chunk_text(text: str, max_tokens: int = 300) -> list[str]:
    lines = [line.strip() for line in text.split("\n") if line.strip()]
    chunks, current = [], []

    for line in lines:
        current.append(line)
        if len(" ".join(current).split()) >= max_tokens:
            chunk = " ".join(current)
            print(f"📦 Chunk created (approx {len(chunk.split())} words):\n{chunk[:300]}...\n")
            chunks.append(chunk)
            current = []

    if current:
        chunk = " ".join(current)
        print(f"📦 Final chunk (approx {len(chunk.split())} words):\n{chunk[:300]}...\n")
        chunks.append(chunk)

    return chunks
